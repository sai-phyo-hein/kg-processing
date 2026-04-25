"""PowerPoint processing module for handling PPTX files."""

import os
import tempfile
from typing import List

try:
    from PIL import Image, ImageDraw, ImageFont
    from pptx import Presentation
except ImportError as e:
    raise ImportError(f"Required dependencies not installed. Please install: {e}") from e


def _create_text_image(text: str, title: str = "") -> Image.Image:
    """Create an image from text content.

    Args:
        text: Text content to render
        title: Optional title for the image

    Returns:
        PIL Image with rendered text
    """
    # Create a white image
    width, height = 800, 600
    img = Image.new("RGB", (width, height), color="white")
    draw = ImageDraw.Draw(img)

    # Try to use a default font
    try:
        font = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf", 16)
        title_font = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf", 20)
    except (IOError, OSError):
        font = ImageFont.load_default()
        title_font = ImageFont.load_default()

    # Draw title if provided
    if title:
        draw.text((10, 10), title, fill="black", font=title_font)

    # Draw text content (simplified - in production, you'd want proper text wrapping)
    y_offset = 50
    for line in text.split("\n"):
        if y_offset < height - 30:
            draw.text((10, y_offset), line, fill="black", font=font)
            y_offset += 25

    return img


def process_pptx(pptx_path: str) -> List[str]:
    """Process PPTX file and convert slides to images.

    Args:
        pptx_path: Path to the PPTX file

    Returns:
        List of base64 encoded slide images

    Raises:
        FileNotFoundError: If PPTX file doesn't exist
        PermissionError: If PPTX file cannot be read
        IOError: If PPTX file cannot be processed
    """
    import base64

    # Validate file exists
    if not os.path.exists(pptx_path):
        raise FileNotFoundError(f"PPTX file not found: {pptx_path}")

    # Validate file is readable
    if not os.access(pptx_path, os.R_OK):
        raise PermissionError(f"Cannot read PPTX file: {pptx_path}")

    base64_images = []

    try:
        # For PPTX, we'll extract slide content
        prs = Presentation(pptx_path)

        for slide_num, slide in enumerate(prs.slides):
            # Extract text from slide
            text_content = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text_content.append(shape.text)

            # Create image from slide content
            slide_text = "\n".join(text_content) if text_content else f"Slide {slide_num + 1}"
            img = _create_text_image(slide_text, f"Slide {slide_num + 1}")

            buffered = tempfile.NamedTemporaryFile(suffix=".png", delete=False)
            img.save(buffered, format="PNG")
            buffered.close()

            with open(buffered.name, "rb") as f:
                image_data = f.read()
                base64_images.append(base64.b64encode(image_data).decode("utf-8"))

            os.unlink(buffered.name)
    except Exception as e:
        raise IOError(f"Failed to process PPTX: {e}") from e

    return base64_images
